import io
import os
import fitz  # PyMuPDF
import requests
import tempfile
import hashlib
import mimetypes
from pathlib import Path
import zipfile
import shutil
from concurrent.futures import ThreadPoolExecutor, as_completed
from typing import Dict, List, Tuple, Optional, Union
import json

from .vision_ocr import extract_text_with_gemini_vision

# Document parsing libraries
try:
    import docx
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

# Import table processor
try:
    from .table_processor import TableProcessor
    TABLE_PROCESSOR_AVAILABLE = True
except ImportError:
    TABLE_PROCESSOR_AVAILABLE = False

try:
    import openpyxl
    from openpyxl import load_workbook
    XLSX_AVAILABLE = True
except ImportError:
    XLSX_AVAILABLE = False

try:
    import pandas as pd
    PANDAS_AVAILABLE = True
except ImportError:
    PANDAS_AVAILABLE = False

try:
    import pytesseract
    from PIL import Image
    OCR_AVAILABLE = True
    # If the executable isn't on PATH, allow env var or fallback path to specify it.
    # F:\ fallback is intentional: Tesseract was installed on F: due to C: drive space constraints.
    if OCR_AVAILABLE:
        custom_tess_cmd = os.environ.get("TESSERACT_CMD") or r"F:\\Tesseract\\tesseract.exe"
        if os.path.exists(custom_tess_cmd):
            pytesseract.pytesseract.tesseract_cmd = custom_tess_cmd
except ImportError:
    OCR_AVAILABLE = False

try:
    import docx2txt
    DOC2TXT_AVAILABLE = True
except ImportError:
    DOC2TXT_AVAILABLE = False

try:
    import xlrd
    XLRD_AVAILABLE = True
except ImportError:
    XLRD_AVAILABLE = False

try:
    import pptx
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


def _format_combined_ocr_output(tesseract_text: str, gemini_text: str) -> str:
    """Primary downstream text is Gemini; Tesseract is appended as supplement."""
    tesseract_text = (tesseract_text or "").strip()
    gemini_text = (gemini_text or "").strip()

    if gemini_text and tesseract_text:
        return f"[Gemini Vision]:\n{gemini_text}\n\n[Tesseract OCR supplement]:\n{tesseract_text}"
    if gemini_text:
        return f"[Gemini Vision]:\n{gemini_text}"
    if tesseract_text:
        return f"[Tesseract OCR]:\n{tesseract_text}"
    return ""


class DocumentParser:
    """Comprehensive document parser supporting multiple formats with OCR capabilities."""
    
    def __init__(self):
        self.supported_extensions = {
            '.pdf': self._parse_pdf,
            '.docx': self._parse_docx,
            '.doc': self._parse_doc,
            '.xlsx': self._parse_xlsx,
            '.xls': self._parse_xls,
            '.net': self._parse_net_file,
            '.cs': self._parse_net_file,
            '.vb': self._parse_net_file,
            '.aspx': self._parse_net_file,
            '.ashx': self._parse_net_file,
            '.asmx': self._parse_net_file,
            '.config': self._parse_net_file,
            '.csproj': self._parse_net_file,
            '.vbproj': self._parse_net_file,
            '.sln': self._parse_net_file,
            '.dll': self._parse_binary_file,
            '.exe': self._parse_binary_file,
            '.pdb': self._parse_binary_file,
            '.pptx': self._parse_pptx,
            '.png': self._parse_image,
            '.jpg': self._parse_image,
            '.jpeg': self._parse_image,
            '.zip': self._parse_zip
        }
        
        # OCR settings
        self.ocr_config = '--oem 3 --psm 6'
        self.ocr_languages = ['eng']
        
    def download_file(self, url: str) -> str:
        """Download file from URL and return local path."""
        response = requests.get(url)
        if response.status_code != 200:
            raise Exception(f"Failed to download file: {response.status_code}")
        
        # Determine file extension from URL or content-type
        content_type = response.headers.get('content-type', '')
        extension = self._get_extension_from_url(url) or self._get_extension_from_mime(content_type)
        
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=extension)
        temp_file.write(response.content)
        temp_file.close()
        return temp_file.name
    
    def get_file_hash(self, file_path: str) -> str:
        """Calculate SHA-256 hash of a file."""
        with open(file_path, "rb") as f:
            file_bytes = f.read()
            sha256_hash = hashlib.sha256(file_bytes).hexdigest()
        return sha256_hash
    
    def extract_text_from_file(self, file_path: str) -> Dict[str, Union[str, List[Dict]]]:
        """Extract text from file based on its extension."""
        file_path = Path(file_path)
        extension = file_path.suffix.lower()
        
        if extension not in self.supported_extensions:
            raise ValueError(f"Unsupported file extension: {extension}")
        
        parser_func = self.supported_extensions[extension]
        return parser_func(file_path)
    
    def _get_extension_from_url(self, url: str) -> Optional[str]:
        """Extract file extension from URL."""
        path = url.split('?')[0]  # Remove query parameters
        return Path(path).suffix.lower()
    
    def _get_extension_from_mime(self, mime_type: str) -> Optional[str]:
        """Get file extension from MIME type."""
        mime_to_ext = {
            'application/pdf': '.pdf',
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document': '.docx',
            'application/msword': '.doc',
            'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': '.xlsx',
            'application/vnd.ms-excel': '.xls',
            'text/plain': '.txt',
            'text/x-csharp': '.cs',
            'text/x-vb': '.vb',
            'application/xml': '.xml'
        }
        return mime_to_ext.get(mime_type)
    
    def _ocr_pdf_image(self, page_num: int, img_index: int, img_data: bytes) -> Tuple[int, int, str]:
        """Run Tesseract + Gemini Vision on a single PDF embedded image."""
        ocr_text = ""
        if OCR_AVAILABLE:
            try:
                img_pil = Image.open(io.BytesIO(img_data))
                ocr_text = pytesseract.image_to_string(img_pil, config=self.ocr_config)
            except Exception as e:
                print(f"Tesseract error on image {img_index + 1} on page {page_num + 1}: {e}")

        gemini_text = extract_text_with_gemini_vision(img_data, ocr_text)
        combined = _format_combined_ocr_output(ocr_text, gemini_text)
        return page_num, img_index, combined

    def _parse_pdf(self, file_path: Path) -> Dict[str, Union[str, List[Dict]]]:
        """Parse PDF file with OCR support for images."""
        text_content = ""
        tables_data = []
        image_jobs: List[Tuple[int, int, bytes]] = []

        with fitz.open(file_path) as doc:
            for page_num, page in enumerate(doc):
                page_text = page.get_text()
                text_content += f"\n--- Page {page_num + 1} ---\n{page_text}\n"

                for img_index, img in enumerate(page.get_images()):
                    try:
                        xref = img[0]
                        pix = fitz.Pixmap(doc, xref)
                        if pix.n - pix.alpha < 4:  # GRAY or RGB
                            img_data = pix.tobytes("png")
                            image_jobs.append((page_num, img_index, img_data))
                        pix = None
                    except Exception as e:
                        print(f"Error extracting image {img_index} on page {page_num + 1}: {e}")

        if image_jobs:
            max_workers = min(8, len(image_jobs))
            ocr_results: List[Tuple[int, int, str]] = []
            with ThreadPoolExecutor(max_workers=max_workers) as executor:
                futures = [
                    executor.submit(self._ocr_pdf_image, page_num, img_index, img_data)
                    for page_num, img_index, img_data in image_jobs
                ]
                for future in as_completed(futures):
                    try:
                        ocr_results.append(future.result())
                    except Exception as e:
                        print(f"Error during parallel PDF image OCR: {e}")

            for page_num, img_index, combined in sorted(ocr_results, key=lambda x: (x[0], x[1])):
                if combined.strip():
                    text_content += (
                        f"\n[OCR Text from Image {img_index + 1} on Page {page_num + 1}]:\n"
                        f"{combined}\n"
                    )

        return {
            "text": text_content,
            "tables": tables_data,
            "file_type": "pdf"
        }
    
    def _parse_docx(self, file_path: Path) -> Dict[str, Union[str, List[Dict]]]:
        """Parse DOCX file with table extraction."""
        if not DOCX_AVAILABLE:
            raise ImportError("python-docx is required for DOCX parsing")
        
        text_content = ""
        tables_data = []
        
        try:
            doc = Document(file_path)
            
            # Extract text and tables
            for element in doc.element.body:
                if element.tag.endswith('p'):  # Paragraph
                    text_content += element.text + "\n"
                elif element.tag.endswith('tbl'):  # Table
                    table_data = self._extract_table_from_docx(element)
                    tables_data.append(table_data)
                    text_content += f"\n[Table {len(tables_data)}]:\n{table_data['text']}\n"
        
        except Exception as e:
            print(f"Error parsing DOCX: {e}")
            # Fallback to docx2txt
            if DOC2TXT_AVAILABLE:
                text_content = docx2txt.process(str(file_path))
        
        return {
            "text": text_content,
            "tables": tables_data,
            "file_type": "docx"
        }
    
    def _parse_doc(self, file_path: Path) -> Dict[str, Union[str, List[Dict]]]:
        """Parse DOC file using docx2txt."""
        if not DOC2TXT_AVAILABLE:
            raise ImportError("docx2txt is required for DOC parsing")
        
        try:
            text_content = docx2txt.process(str(file_path))
        except Exception as e:
            print(f"Error parsing DOC: {e}")
            text_content = f"Error parsing DOC file: {e}"
        
        return {
            "text": text_content,
            "tables": [],
            "file_type": "doc"
        }
    
    def _parse_xlsx(self, file_path: Path) -> Dict[str, Union[str, List[Dict]]]:
        """Parse XLSX file with comprehensive table extraction."""
        if not XLSX_AVAILABLE:
            raise ImportError("openpyxl is required for XLSX parsing")
        
        text_content = ""
        tables_data = []
        
        try:
            # Use specialized table processor if available
            if TABLE_PROCESSOR_AVAILABLE:
                processor = TableProcessor()
                processed_tables = processor.process_xlsx_tables(str(file_path))
                
                if "error" not in processed_tables:
                    # Add summary
                    summary = processor.get_table_summary(processed_tables["sheets"])
                    text_content += f"XLSX File Analysis:\n{summary}\n\n"
                    
                    # Add detailed table information
                    for sheet_name, tables in processed_tables["sheets"].items():
                        text_content += f"\n--- Sheet: {sheet_name} ---\n"
                        for table in tables:
                            text_content += f"\nTable {table['table_index']}:\n"
                            text_content += table['text'] + "\n"
                            tables_data.append(table)
                else:
                    # Fallback to basic parsing
                    text_content = self._basic_xlsx_parsing(file_path)
            else:
                # Fallback to basic parsing
                text_content = self._basic_xlsx_parsing(file_path)
        
        except Exception as e:
            print(f"Error parsing XLSX: {e}")
            text_content = f"Error parsing XLSX file: {e}"
        
        return {
            "text": text_content,
            "tables": tables_data,
            "file_type": "xlsx"
        }
    
    def _basic_xlsx_parsing(self, file_path: Path) -> str:
        """Basic XLSX parsing fallback."""
        text_content = ""
        
        try:
            workbook = load_workbook(file_path, data_only=True)
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text_content += f"\n--- Sheet: {sheet_name} ---\n"
                
                # Extract table data
                table_data = self._extract_table_from_xlsx(sheet, sheet_name)
                
                # Add table text to main content
                text_content += table_data['text'] + "\n"
        
        except Exception as e:
            print(f"Error in basic XLSX parsing: {e}")
            text_content = f"Error parsing XLSX file: {e}"
        
        return text_content
    
    def _parse_xls(self, file_path: Path) -> Dict[str, Union[str, List[Dict]]]:
        """Parse XLS file using xlrd."""
        if not XLRD_AVAILABLE:
            raise ImportError("xlrd is required for XLS parsing")
        
        text_content = ""
        tables_data = []
        
        try:
            workbook = xlrd.open_workbook(file_path)
            
            for sheet_name in workbook.sheet_names():
                sheet = workbook.sheet_by_name(sheet_name)
                text_content += f"\n--- Sheet: {sheet_name} ---\n"
                
                # Extract table data
                table_data = self._extract_table_from_xls(sheet, sheet_name)
                tables_data.append(table_data)
                
                # Add table text to main content
                text_content += table_data['text'] + "\n"
        
        except Exception as e:
            print(f"Error parsing XLS: {e}")
            text_content = f"Error parsing XLS file: {e}"
        
        return {
            "text": text_content,
            "tables": tables_data,
            "file_type": "xls"
        }
    
    def _parse_net_file(self, file_path: Path) -> Dict[str, Union[str, List[Dict]]]:
        """Parse .NET related files (C#, VB, config files, etc.)."""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
        except UnicodeDecodeError:
            try:
                with open(file_path, 'r', encoding='latin-1', errors='ignore') as f:
                    content = f.read()
            except Exception as e:
                content = f"Error reading .NET file: {e}"
        
        # Add file metadata
        file_info = f"File: {file_path.name}\nType: .NET File\nExtension: {file_path.suffix}\n\n"
        text_content = file_info + content
        
        return {
            "text": text_content,
            "tables": [],
            "file_type": "net"
        }
    
    def _parse_binary_file(self, file_path: Path) -> Dict[str, Union[str, List[Dict]]]:
        """Parse binary files (DLL, EXE, PDB) - extract basic info only."""
        try:
            with open(file_path, 'rb') as f:
                # Read first 1024 bytes for basic analysis
                header = f.read(1024)
            
            # Basic binary file analysis
            file_size = file_path.stat().st_size
            text_content = f"""Binary File Analysis:
File: {file_path.name}
Size: {file_size} bytes
Type: {file_path.suffix.upper()} file
Header (hex): {header[:64].hex()}

Note: This is a binary file. For detailed analysis, use specialized tools.
"""
        except Exception as e:
            text_content = f"Error analyzing binary file: {e}"
        
        return {
            "tables": [],
            "file_type": "binary"
        }
    
    def _parse_image(self, file_path: Path) -> Dict[str, Union[str, List[Dict]]]:
        """Parse an image file using OCR to extract text."""
        try:
            with open(file_path, "rb") as f:
                image_bytes = f.read()
        except Exception as e:
            return {
                "text": f"Error reading image file: {e}",
                "tables": [],
                "file_type": "image",
            }

        tesseract_text = ""
        if OCR_AVAILABLE:
            try:
                from PIL import Image as PILImage
                img = PILImage.open(io.BytesIO(image_bytes))
                tesseract_text = pytesseract.image_to_string(img, config=self.ocr_config)
                if tesseract_text.strip():
                    print(f"[OCR Image Tesseract] {tesseract_text.strip()[:200]}")
            except Exception as e:
                print(f"Tesseract error parsing image file: {e}")
        else:
            print(f"[OCR Image] Tesseract not available for {file_path.name}; using Gemini Vision only")

        gemini_text = extract_text_with_gemini_vision(image_bytes, tesseract_text)
        text = _format_combined_ocr_output(tesseract_text, gemini_text)

        if not text.strip():
            text = f"[Image file {file_path.name} – no text extracted, {len(image_bytes)} bytes]"

        if text.strip():
            print(f"[OCR Image combined] {text.strip()[:200]}")

        return {
            "text": text,
            "tables": [],
            "file_type": "image",
        }

    def _parse_pptx(self, file_path: Path) -> Dict[str, Union[str, List[Dict]]]:
        """Parse PPTX file with slide-level parallelism, extracting text and images,
        and performing OCR on images when available. The resulting OCR text is
        appended to the slide text so it can be chunked and embedded later."""
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx is required for PPTX parsing")

        import base64
        from pptx import Presentation

        text_chunks: list[str] = []
        images_data: list[Dict[str, Union[int, str]]] = []

        def _ocr_slide_image(slide_index: int, image_bytes: bytes) -> str:
            ocr_text = ""
            if OCR_AVAILABLE:
                try:
                    from PIL import Image as PILImage
                    ocr_img = PILImage.open(io.BytesIO(image_bytes))
                    ocr_text = pytesseract.image_to_string(ocr_img, config=self.ocr_config)
                except Exception as ocr_err:
                    print(f"Tesseract OCR error on slide {slide_index + 1}: {ocr_err}")

            gemini_text = extract_text_with_gemini_vision(image_bytes, ocr_text)
            combined = _format_combined_ocr_output(ocr_text, gemini_text)
            if combined.strip():
                print(f"[OCR Slide {slide_index + 1}] {combined.strip()[:200]}")
            return combined

        def _process_slide(slide_index: int, slide) -> Tuple[str, list[Dict]]:
            """Extract text, images and OCR text from a single slide."""
            slide_lines: list[str] = [f"\n--- Slide {slide_index + 1} ---"]
            local_images: list[Dict] = []
            slide_image_bytes: list[bytes] = []

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_lines.append(shape.text)

                try:
                    if getattr(shape, "shape_type", None) == 13 and hasattr(shape, "image"):
                        image = shape.image
                        image_bytes = image.blob
                        slide_image_bytes.append(image_bytes)

                        image_b64 = base64.b64encode(image_bytes).decode("utf-8")
                        local_images.append({
                            "slide": slide_index + 1,
                            "image_format": image.ext,
                            "image_base64": image_b64
                        })
                        slide_lines.append(
                            f"[Image on Slide {slide_index + 1}: {image.ext}, {len(image_bytes)} bytes]"
                        )
                except Exception as img_err:
                    print(f"Error extracting image on slide {slide_index + 1}: {img_err}")

            if slide_image_bytes:
                if len(slide_image_bytes) == 1:
                    combined = _ocr_slide_image(slide_index, slide_image_bytes[0])
                    if combined.strip():
                        slide_lines.append(f"[OCR] {combined.strip()}")
                else:
                    max_workers = min(8, len(slide_image_bytes))
                    with ThreadPoolExecutor(max_workers=max_workers) as executor:
                        futures = [
                            executor.submit(_ocr_slide_image, slide_index, image_bytes)
                            for image_bytes in slide_image_bytes
                        ]
                        for future in as_completed(futures):
                            try:
                                combined = future.result()
                                if combined.strip():
                                    slide_lines.append(f"[OCR] {combined.strip()}")
                            except Exception as ocr_err:
                                print(f"OCR error on slide {slide_index + 1}: {ocr_err}")

            return "\n".join(slide_lines), local_images

        prs = Presentation(file_path)

        # Use up to 8 threads or number of slides, whichever is smaller
        max_workers = min(8, len(prs.slides)) or 1
        with ThreadPoolExecutor(max_workers=max_workers) as executor:
            futures = {executor.submit(_process_slide, idx, slide): idx for idx, slide in enumerate(prs.slides)}
            for future in as_completed(futures):
                slide_text, slide_images = future.result()
                text_chunks.append(slide_text)
                images_data.extend(slide_images)

        # Preserve slide order by sorting on the marker we added (slide number)
        text_content = "\n".join(sorted(text_chunks, key=lambda x: int(x.split("--- Slide ")[1].split(" ")[0])))

        return {
            "text": text_content,
            "images": images_data,
            "tables": [],
            "file_type": "pptx"
        }
    
    def _extract_table_from_docx(self, table_element) -> Dict:
        """Extract table data from DOCX table element."""
        table_data = []
        table_text = ""
        
        try:
            for row in table_element.findall('.//w:tr', {'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'}):
                row_data = []
                for cell in row.findall('.//w:tc', {'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'}):
                    cell_text = ""
                    for text_elem in cell.findall('.//w:t', {'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'}):
                        cell_text += text_elem.text or ""
                    row_data.append(cell_text.strip())
                    table_text += f"{cell_text.strip()}\t"
                table_data.append(row_data)
                table_text += "\n"
        except Exception as e:
            print(f"Error extracting table from DOCX: {e}")
        
        return {
            "data": table_data,
            "text": table_text,
            "rows": len(table_data),
            "columns": len(table_data[0]) if table_data else 0
        }
    
    def _extract_table_from_xlsx(self, sheet, sheet_name: str) -> Dict:
        """Extract comprehensive table data from XLSX sheet."""
        table_data = []
        table_text = f"Sheet: {sheet_name}\n"
        
        try:
            # Get the used range
            max_row = sheet.max_row
            max_col = sheet.max_column
            
            # Extract all data
            for row in range(1, max_row + 1):
                row_data = []
                for col in range(1, max_col + 1):
                    cell = sheet.cell(row=row, column=col)
                    cell_value = cell.value
                    
                    # Handle different data types
                    if cell_value is None:
                        cell_value = ""
                    elif isinstance(cell_value, (int, float)):
                        cell_value = str(cell_value)
                    elif isinstance(cell_value, str):
                        cell_value = cell_value.strip()
                    
                    row_data.append(cell_value)
                    table_text += f"{cell_value}\t"
                
                table_data.append(row_data)
                table_text += "\n"
        
        except Exception as e:
            print(f"Error extracting table from XLSX: {e}")
        
        return {
            "data": table_data,
            "text": table_text,
            "rows": len(table_data),
            "columns": len(table_data[0]) if table_data else 0,
            "sheet_name": sheet_name
        }
    
    def _extract_table_from_xls(self, sheet, sheet_name: str) -> Dict:
        """Extract table data from XLS sheet."""
        table_data = []
        table_text = f"Sheet: {sheet_name}\n"
        
        try:
            for row_idx in range(sheet.nrows):
                row_data = []
                for col_idx in range(sheet.ncols):
                    cell_value = sheet.cell_value(row_idx, col_idx)
                    
                    # Handle different data types
                    if cell_value is None:
                        cell_value = ""
                    elif isinstance(cell_value, (int, float)):
                        cell_value = str(cell_value)
                    elif isinstance(cell_value, str):
                        cell_value = cell_value.strip()
                    
                    row_data.append(cell_value)
                    table_text += f"{cell_value}\t"
                
                table_data.append(row_data)
                table_text += "\n"
        
        except Exception as e:
            print(f"Error extracting table from XLS: {e}")
        
        return {
            "data": table_data,
            "text": table_text,
            "rows": len(table_data),
            "columns": len(table_data[0]) if table_data else 0,
            "sheet_name": sheet_name
        }
    
    def extract_tables_from_content(self, content: Dict) -> List[Dict]:
        """Extract and format tables from parsed content."""
        tables = content.get("tables", [])
        formatted_tables = []
        
        for i, table in enumerate(tables):
            formatted_table = {
                "table_index": i,
                "rows": table.get("rows", 0),
                "columns": table.get("columns", 0),
                "data": table.get("data", []),
                "text": table.get("text", ""),
                "sheet_name": table.get("sheet_name", "")
            }
            formatted_tables.append(formatted_table)
        
        return formatted_tables


    def _parse_zip(self, file_path: Path) -> Dict[str, Union[str, List[Dict]]]:
        """Parse ZIP file by extracting its contents (up to two levels) and processing contained documents.
        If after two extraction rounds no non-zip files are found, return an error."""
        temp_dir = tempfile.mkdtemp(prefix="zip_extract_")
        collected_text = ""
        processed_files = []

        def _extract(zip_path: Path, dest: str):
            try:
                with zipfile.ZipFile(zip_path, 'r') as zf:
                    zf.extractall(dest)
                return True
            except Exception as exc:
                print(f"Error extracting {zip_path.name}: {exc}")
                return False

        # First extraction
        level1_dir = os.path.join(temp_dir, "level1")
        os.makedirs(level1_dir, exist_ok=True)
        if not _extract(file_path, level1_dir):
            shutil.rmtree(temp_dir)
            return {"text": f"Error extracting ZIP file {file_path.name}", "tables": [], "file_type": "zip"}

        # Gather files
        level1_files = [Path(os.path.join(root, f)) for root, _, files in os.walk(level1_dir) for f in files]
        non_zip_files = [p for p in level1_files if p.suffix.lower() != '.zip']

        # If still all zips, extract one more level
        if not non_zip_files:
            level2_dir = os.path.join(temp_dir, "level2")
            os.makedirs(level2_dir, exist_ok=True)
            for z in level1_files:
                _extract(z, level2_dir)
            non_zip_files = [Path(os.path.join(root, f)) for root, _, files in os.walk(level2_dir) for f in files if Path(f).suffix.lower() != '.zip']

        if not non_zip_files:
            shutil.rmtree(temp_dir)
            return {"text": "Error: ZIP contains only nested ZIP files and no supported documents.", "tables": [], "file_type": "zip"}

        # Process each non-zip file using existing parser
        parser = DocumentParser()
        for p in non_zip_files:
            try:
                result = parser.extract_text_from_file(str(p))
                collected_text += f"\n--- File: {p.name} ---\n" + result.get("text", "") + "\n"
                processed_files.append(p.name)
            except Exception as err:
                print(f"Error parsing file {p}: {err}")

        shutil.rmtree(temp_dir)
        return {"text": collected_text or "", "tables": [], "file_type": "zip", "processed_files": processed_files}

# Backward compatibility functions
def download_pdf(url: str) -> str:
    """Download PDF file from URL (backward compatibility)."""
    parser = DocumentParser()
    return parser.download_file(url)

def get_file_hash(file_path: str) -> str:
    """Calculate file hash (backward compatibility)."""
    parser = DocumentParser()
    return parser.get_file_hash(file_path)

def extract_text_from_pdf(pdf_path: str) -> str:
    """Extract text from PDF (backward compatibility)."""
    parser = DocumentParser()
    result = parser.extract_text_from_file(pdf_path)
    return result["text"] 